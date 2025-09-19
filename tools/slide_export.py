# tools/slide_export.py
import os
import datetime
from typing import List, Dict, Any
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN

# Basic corporate-ish theme choices
TITLE_COLOR = "203864"      # dark blue
ACCENT_COLOR = "10A37F"     # green accent
BODY_COLOR = "2E2E2E"       # neutral dark
LIGHT_GREY = "777777"

def _add_title_slide(prs: Presentation, title: str, subtitle: str = ""):
    slide_layout = prs.slide_layouts[0]  # Title slide
    slide = prs.slides.add_slide(slide_layout)
    title_tf = slide.shapes.title
    subtitle_tf = slide.placeholders[1]

    title_tf.text = title
    title_tf.text_frame.paragraphs[0].font.size = Pt(42)
    title_tf.text_frame.paragraphs[0].font.bold = True
    title_tf.text_frame.paragraphs[0].font.color.rgb = _hex_to_rgb(TITLE_COLOR)

    subtitle_tf.text = subtitle
    p = subtitle_tf.text_frame.paragraphs[0]
    p.font.size = Pt(16)
    p.font.color.rgb = _hex_to_rgb(LIGHT_GREY)
    p.alignment = PP_ALIGN.LEFT

def _add_section_title_slide(prs: Presentation, title: str, subtitle: str = ""):
    slide_layout = prs.slide_layouts[5]  # Title Only
    slide = prs.slides.add_slide(slide_layout)
    title_tf = slide.shapes.title
    title_tf.text = title
    title_tf.text_frame.paragraphs[0].font.size = Pt(32)
    title_tf.text_frame.paragraphs[0].font.bold = True
    title_tf.text_frame.paragraphs[0].font.color.rgb = _hex_to_rgb(TITLE_COLOR)

    if subtitle:
        txbox = slide.shapes.add_textbox(Inches(1), Inches(1.6), Inches(8), Inches(1))
        tf = txbox.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = subtitle
        p.font.size = Pt(16)
        p.font.color.rgb = _hex_to_rgb(LIGHT_GREY)

def _add_bullets_slide(prs: Presentation, title: str, bullets: List[str]):
    slide_layout = prs.slide_layouts[1]  # Title and Content
    slide = prs.slides.add_slide(slide_layout)
    title_tf = slide.shapes.title
    body_tf = slide.placeholders[1]

    title_tf.text = title
    title_tf.text_frame.paragraphs[0].font.size = Pt(28)
    title_tf.text_frame.paragraphs[0].font.bold = True
    title_tf.text_frame.paragraphs[0].font.color.rgb = _hex_to_rgb(TITLE_COLOR)

    tf = body_tf.text_frame
    tf.clear()
    for i, bullet in enumerate(bullets):
        p = tf.add_paragraph() if i > 0 else tf.paragraphs[0]
        p.text = bullet
        p.level = 0
        p.font.size = Pt(16)
        p.font.color.rgb = _hex_to_rgb(BODY_COLOR)

def _add_step_slide(
    prs: Presentation,
    step_title: str,
    user_input: str,
    goals: List[str],
    best_practices: List[str],
):
    # Title + two-column layout (title-only + custom textboxes)
    slide_layout = prs.slide_layouts[5]  # Title Only
    slide = prs.slides.add_slide(slide_layout)
    slide.shapes.title.text = step_title
    title_p = slide.shapes.title.text_frame.paragraphs[0]
    title_p.font.size = Pt(28)
    title_p.font.bold = True
    title_p.font.color.rgb = _hex_to_rgb(TITLE_COLOR)

    # Left column: user's input / outcome
    left = slide.shapes.add_textbox(Inches(0.75), Inches(1.6), Inches(4.3), Inches(4.5))
    tf_left = left.text_frame
    tf_left.word_wrap = True

    p0 = tf_left.paragraphs[0]
    p0.text = "Your Input"
    p0.font.size = Pt(16)
    p0.font.bold = True
    p0.font.color.rgb = _hex_to_rgb(ACCENT_COLOR)

    p1 = tf_left.add_paragraph()
    p1.text = user_input.strip() or "—"
    p1.font.size = Pt(15)
    p1.font.color.rgb = _hex_to_rgb(BODY_COLOR)

    # Right column: goals + best practices
    right = slide.shapes.add_textbox(Inches(5.0), Inches(1.6), Inches(4.3), Inches(4.5))
    tf_right = right.text_frame
    tf_right.word_wrap = True

    pr0 = tf_right.paragraphs[0]
    pr0.text = "Goals"
    pr0.font.size = Pt(16)
    pr0.font.bold = True
    pr0.font.color.rgb = _hex_to_rgb(ACCENT_COLOR)

    if goals:
        for g in goals:
            pg = tf_right.add_paragraph()
            pg.text = f"• {g}"
            pg.font.size = Pt(14)
            pg.font.color.rgb = _hex_to_rgb(BODY_COLOR)
    else:
        pg = tf_right.add_paragraph()
        pg.text = "—"
        pg.font.size = Pt(14)
        pg.font.color.rgb = _hex_to_rgb(BODY_COLOR)

    pr1 = tf_right.add_paragraph()
    pr1.text = ""  # spacer

    pr2 = tf_right.add_paragraph()
    pr2.text = "Best Practices"
    pr2.font.size = Pt(16)
    pr2.font.bold = True
    pr2.font.color.rgb = _hex_to_rgb(ACCENT_COLOR)

    if best_practices:
        for bp in best_practices:
            pb = tf_right.add_paragraph()
            pb.text = f"• {bp}"
            pb.font.size = Pt(14)
            pb.font.color.rgb = _hex_to_rgb(BODY_COLOR)
    else:
        pb = tf_right.add_paragraph()
        pb.text = "—"
        pb.font.size = Pt(14)
        pb.font.color.rgb = _hex_to_rgb(BODY_COLOR)

def _hex_to_rgb(hex_str: str):
    from pptx.dml.color import RGBColor
    hex_str = hex_str.lstrip("#")
    return RGBColor(int(hex_str[0:2], 16), int(hex_str[2:4], 16), int(hex_str[4:6], 16))

def export_dmaic_to_pptx(steps_filled: List[Dict[str, Any]]) -> str:
    """
    Create a .pptx summarizing a DMAIC session.

    steps_filled: list of dicts with keys:
      - step: str (e.g., 'Define', 'Measure', ...)
      - user_input: str
      - goals: List[str]
      - best_practices: List[str]

    Returns: file path to the generated pptx.
    """
    prs = Presentation()

    # Title slide
    _add_title_slide(
        prs,
        title="DMAIC Summary",
        subtitle=f"Generated {datetime.datetime.now().strftime('%Y-%m-%d %H:%M')}"
    )

    # Overview slide
    _add_section_title_slide(prs, "Overview", subtitle="Your DMAIC session at a glance")

    # Quick index bullets
    index_bullets = [f"{i+1}. {item.get('step','Step')}" for i, item in enumerate(steps_filled)]
    _add_bullets_slide(prs, "Steps Covered", index_bullets)

    # One slide per step
    for item in steps_filled:
        _add_step_slide(
            prs,
            step_title=item.get("step", "Step"),
            user_input=item.get("user_input", ""),
            goals=item.get("goals", []),
            best_practices=item.get("best_practices", []),
        )

    # Save to disk (tmp folder by default)
    os.makedirs("exports", exist_ok=True)
    out_path = os.path.join("exports", f"dmaic_summary_{datetime.datetime.now().strftime('%Y%m%d_%H%M%S')}.pptx")
    prs.save(out_path)
    return out_path
